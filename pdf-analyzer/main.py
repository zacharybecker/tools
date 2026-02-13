import os
import base64
import asyncio
import logging
import tempfile
import subprocess
import uuid
from io import BytesIO
from contextlib import asynccontextmanager

import fitz
import httpx
from fastapi import FastAPI, HTTPException
from pydantic import BaseModel, validator
from tenacity import retry, stop_after_attempt, wait_exponential, retry_if_exception
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from PIL import Image

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

LITELLM_BASE_URL = os.getenv("LITELLM_BASE_URL")
LITELLM_MODEL = os.getenv("LITELLM_MODEL")
MAX_CONCURRENT = int(os.getenv("MAX_CONCURRENT_LLM_CALLS", "5"))
MAX_FILE_SIZE_MB = int(os.getenv("MAX_FILE_SIZE_MB", "50"))
DPI = 300
MAX_DIMENSION_PX = int(os.getenv("MAX_DIMENSION_PX", "3000"))
OVERLAP_PX = 100
JPG_QUALITY = int(os.getenv("JPG_QUALITY", "95"))


def should_retry(exception):
    if isinstance(exception, httpx.HTTPStatusError):
        return exception.response.status_code in [429, 500, 502, 503, 504]
    return isinstance(exception, (httpx.TimeoutException, httpx.ConnectError))


class ConvertRequest(BaseModel):
    pdf_data: str

    @validator('pdf_data')
    def validate_base64(cls, v):
        try:
            base64.b64decode(v, validate=True)
            return v
        except:
            raise ValueError("Invalid base64")


class ConvertPresentationRequest(BaseModel):
    presentation_data: str
    format: str = "pptx"

    @validator('format')
    def validate_format(cls, v):
        if v not in ("ppt", "pptx"):
            raise ValueError("Format must be 'ppt' or 'pptx'")
        return v

    @validator('presentation_data')
    def validate_base64(cls, v):
        try:
            base64.b64decode(v, validate=True)
            return v
        except:
            raise ValueError("Invalid base64")


class ConvertResponse(BaseModel):
    markdown: str
    total_pages: int
    successful_pages: int
    failed_pages: int


class ConvertPresentationToPdfResponse(BaseModel):
    pdf_data: str
    pages: int


@asynccontextmanager
async def lifespan(app: FastAPI):
    app.state.client = httpx.AsyncClient(timeout=120.0)
    app.state.semaphore = asyncio.Semaphore(MAX_CONCURRENT)
    yield
    await app.state.client.aclose()


app = FastAPI(lifespan=lifespan)


@retry(
    stop=stop_after_attempt(10),
    wait=wait_exponential(multiplier=2, min=4, max=120),
    retry=retry_if_exception(should_retry)
)
async def call_litellm(client, semaphore, image_bytes, label,
                       prompt="Convert this PDF page section to markdown. Preserve structure and formatting."):
    async with semaphore:
        image_base64 = base64.b64encode(image_bytes).decode()
        response = await client.post(
            f"{LITELLM_BASE_URL}/chat/completions",
            json={
                "model": LITELLM_MODEL,
                "messages": [{
                    "role": "user",
                    "content": [
                        {"type": "text", "text": prompt},
                        {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{image_base64}"}}
                    ]
                }],
                "max_tokens": 4096
            }
        )

        if response.status_code == 429:
            retry_after = int(response.headers.get("Retry-After", "30"))
            logger.warning(f"Rate limited on {label}, waiting {retry_after}s before retry")
            await asyncio.sleep(retry_after)

        response.raise_for_status()
        return response.json()["choices"][0]["message"]["content"]


async def convert_chunk(client, semaphore, image_bytes, page_num, chunk_num, total_chunks):
    try:
        label = f"page {page_num + 1}" if total_chunks == 1 else f"page {page_num + 1} chunk {chunk_num + 1}/{total_chunks}"
        markdown = await call_litellm(client, semaphore, image_bytes, label)
        return {"success": True, "page": page_num, "chunk": chunk_num, "markdown": markdown}
    except Exception as e:
        logger.error(f"Page {page_num + 1} chunk {chunk_num + 1} failed: {e}")
        return {"success": False, "page": page_num, "chunk": chunk_num, "error": str(e)}


def render_pdf_to_chunks(pdf_bytes):
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    zoom = DPI / 72
    matrix = fitz.Matrix(zoom, zoom)

    chunks = []
    for page_num, page in enumerate(doc):
        width_px = page.rect.width * zoom
        height_px = page.rect.height * zoom

        if height_px <= MAX_DIMENSION_PX and width_px <= MAX_DIMENSION_PX:
            pix = page.get_pixmap(matrix=matrix)
            chunks.append({"image": pix.tobytes("jpeg", jpg_quality=JPG_QUALITY), "page": page_num, "chunk": 0, "total_chunks": 1})
        else:
            num_chunks = int((max(height_px, width_px) / MAX_DIMENSION_PX)) + 1
            chunk_height_pt = page.rect.height / num_chunks
            overlap_pt = OVERLAP_PX / zoom

            for chunk_idx in range(num_chunks):
                y0 = max(0, chunk_idx * chunk_height_pt - overlap_pt)
                y1 = min(page.rect.height, (chunk_idx + 1) * chunk_height_pt + overlap_pt)
                clip = fitz.Rect(0, y0, page.rect.width, y1)
                pix = page.get_pixmap(matrix=matrix, clip=clip)
                chunks.append({"image": pix.tobytes("jpeg", jpg_quality=JPG_QUALITY), "page": page_num, "chunk": chunk_idx, "total_chunks": num_chunks})

    doc.close()
    return chunks


async def process_chunks(chunks):
    tasks = [convert_chunk(app.state.client, app.state.semaphore, c["image"], c["page"], c["chunk"], c["total_chunks"])
             for c in chunks]
    results = await asyncio.gather(*tasks, return_exceptions=True)

    page_markdowns = {}
    for r in results:
        if isinstance(r, dict) and r.get("success"):
            page_num = r["page"]
            if page_num not in page_markdowns:
                page_markdowns[page_num] = []
            page_markdowns[page_num].append((r["chunk"], r["markdown"]))

    total_pages = len(set(c["page"] for c in chunks))
    successful_pages = len(page_markdowns)

    markdown_parts = []
    for page_num in sorted(page_markdowns.keys()):
        chunks_md = [md for _, md in sorted(page_markdowns[page_num])]
        combined_page_md = "\n\n".join(chunks_md)
        markdown_parts.append(f"# Page {page_num + 1}\n\n{combined_page_md}")

    combined = "\n\n---\n\n".join(markdown_parts) if markdown_parts else "# Conversion Failed"

    return ConvertResponse(
        markdown=combined,
        total_pages=total_pages,
        successful_pages=successful_pages,
        failed_pages=total_pages - successful_pages
    )


def apply_inline_formatting(para):
    parts = []
    for run in para.runs:
        text = run.text
        if not text:
            continue
        if run.font.bold and run.font.italic:
            text = f"***{text}***"
        elif run.font.bold:
            text = f"**{text}**"
        elif run.font.italic:
            text = f"*{text}*"
        parts.append(text)
    return "".join(parts)


def extract_text_frame(shape):
    lines = []
    for para in shape.text_frame.paragraphs:
        text = apply_inline_formatting(para)
        if not text.strip():
            lines.append("")
            continue

        # Determine heading level from placeholder type
        if shape.is_placeholder:
            ph_type = shape.placeholder_format.type
            if ph_type == PP_PLACEHOLDER.TITLE or ph_type == PP_PLACEHOLDER.CENTER_TITLE:
                lines.append(f"## {text}")
                continue
            elif ph_type == PP_PLACEHOLDER.SUBTITLE:
                lines.append(f"### {text}")
                continue

        # Bullet indentation from paragraph level
        indent = "  " * para.level
        if para.level > 0:
            lines.append(f"{indent}- {text}")
        else:
            lines.append(text)

    return "\n".join(lines)


def extract_table(table):
    rows = []
    for row in table.rows:
        cells = []
        for cell in row.cells:
            cell_text = cell.text.replace("|", "\\|").replace("\n", " ").strip()
            cells.append(cell_text)
        rows.append("| " + " | ".join(cells) + " |")

    if not rows:
        return ""

    # Insert header separator after first row
    num_cols = len(table.rows[0].cells)
    separator = "| " + " | ".join(["---"] * num_cols) + " |"
    rows.insert(1, separator)

    return "\n".join(rows)


def extract_picture(shape, slide_idx, pos):
    placeholder = f"[IMG_{slide_idx}_{pos}]"
    try:
        image = shape.image
        image_bytes = image.blob
        content_type = image.content_type
        return placeholder, (image_bytes, content_type)
    except Exception as e:
        logger.warning(f"Failed to extract image from slide {slide_idx + 1}: {e}")
        return f"*[Image: could not extract]*", None


def extract_chart(chart):
    parts = []
    if chart.chart_title and chart.chart_title.has_text_frame:
        title = chart.chart_title.text_frame.text
        if title:
            parts.append(f"**Chart: {title}**")

    for series in chart.series:
        if hasattr(series, 'tx') and series.tx and hasattr(series.tx, 'strRef'):
            # Try to get series name
            pass

    try:
        plot = chart.plots[0]
        if hasattr(plot, 'categories') and plot.categories:
            cats = [str(c) for c in plot.categories]
            if cats:
                parts.append(f"Categories: {', '.join(cats)}")
    except Exception:
        pass

    try:
        series_names = []
        for series in chart.series:
            if hasattr(series, 'tx') and series.tx:
                if hasattr(series.tx, 'strRef') and series.tx.strRef:
                    series_names.append(str(series.tx.strRef))
        if series_names:
            parts.append(f"Series: {', '.join(series_names)}")
    except Exception:
        pass

    return "\n".join(parts) if parts else "*[Chart]*"


def extract_group_shape(group, slide_idx, pos):
    parts = []
    image_tasks = []
    for i, shape in enumerate(group.shapes):
        md, imgs = extract_shape(shape, slide_idx, pos * 100 + i)
        if md:
            parts.append(md)
        image_tasks.extend(imgs)
    return "\n".join(parts), image_tasks


def extract_shape(shape, slide_idx, pos):
    image_tasks = []

    if shape.has_table:
        return extract_table(shape.table), image_tasks

    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        md, img = extract_picture(shape, slide_idx, pos)
        if img:
            image_tasks.append((f"[IMG_{slide_idx}_{pos}]", img[0], img[1]))
        return md, image_tasks

    if shape.has_chart:
        return extract_chart(shape.chart), image_tasks

    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        return extract_group_shape(shape, slide_idx, pos)

    if shape.has_text_frame:
        return extract_text_frame(shape), image_tasks

    return "", image_tasks


def extract_slide(slide, slide_idx):
    # Sort shapes by position (top, then left)
    shapes = sorted(slide.shapes, key=lambda s: (s.top or 0, s.left or 0))

    parts = []
    image_tasks = []
    for pos, shape in enumerate(shapes):
        md, imgs = extract_shape(shape, slide_idx, pos)
        if md.strip():
            parts.append(md)
        image_tasks.extend(imgs)

    return "\n\n".join(parts), image_tasks


def normalize_image_to_jpeg(image_bytes, content_type):
    try:
        img = Image.open(BytesIO(image_bytes))
        if img.mode in ("RGBA", "P", "LA"):
            img = img.convert("RGB")
        buf = BytesIO()
        img.save(buf, format="JPEG", quality=JPG_QUALITY)
        return buf.getvalue()
    except Exception as e:
        logger.warning(f"Could not convert image (type={content_type}): {e}")
        return None


async def describe_images_batch(image_tasks):
    if not image_tasks:
        return {}

    client = app.state.client
    semaphore = app.state.semaphore
    prompt = "Describe this image concisely for a markdown document. Focus on the key content and data shown."

    async def describe_one(placeholder, image_bytes, content_type):
        jpeg_bytes = normalize_image_to_jpeg(image_bytes, content_type)
        if jpeg_bytes is None:
            return placeholder, "*[Image: format not supported]*"
        try:
            description = await call_litellm(client, semaphore, jpeg_bytes, f"image {placeholder}", prompt=prompt)
            return placeholder, description
        except Exception as e:
            logger.error(f"Failed to describe image {placeholder}: {e}")
            return placeholder, "*[Image: description unavailable]*"

    results = await asyncio.gather(*[describe_one(ph, ib, ct) for ph, ib, ct in image_tasks])
    return dict(results)


def resolve_image_placeholders(md, descriptions):
    for placeholder, description in descriptions.items():
        md = md.replace(placeholder, description)
    return md


_NS_A = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
_NS_P = '{http://schemas.openxmlformats.org/presentationml/2006/main}'
_SHAPE_TAG_SUFFIXES = ('}sp', '}pic', '}grpSp')


def _remove_decorative_shapes(sp_tree):
    """Remove shapes from an spTree that have no text, images, or placeholder role."""
    for sp in list(sp_tree):
        if not any(sp.tag.endswith(s) for s in _SHAPE_TAG_SUFFIXES):
            continue
        has_text = sp.findall(f'.//{_NS_A}t')
        has_image = sp.findall(f'.//{_NS_A}blip') or sp.tag.endswith('}pic')
        is_placeholder = sp.findall(f'.//{_NS_P}ph')
        if not has_text and not has_image and not is_placeholder:
            sp_tree.remove(sp)


def strip_presentation_backgrounds(file_bytes):
    """Remove backgrounds and decorative shapes from a PPTX presentation."""
    prs = Presentation(BytesIO(file_bytes))

    for slide in prs.slides:
        slide.background.fill.background()

    for layout in prs.slide_layouts:
        layout.background.fill.background()
        _remove_decorative_shapes(layout.shapes._spTree)

    for master in prs.slide_masters:
        master.background.fill.background()
        _remove_decorative_shapes(master.shapes._spTree)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _libreoffice_convert(file_bytes, input_ext, output_fmt):
    """Run LibreOffice headless to convert a file. Returns output file bytes."""
    with tempfile.TemporaryDirectory() as tmpdir:
        input_path = os.path.join(tmpdir, f"input.{input_ext}")
        with open(input_path, "wb") as f:
            f.write(file_bytes)

        profile_dir = os.path.join(tmpdir, f"profile-{uuid.uuid4().hex}")
        result = subprocess.run(
            [
                "libreoffice", "--headless", "--norestore",
                f"-env:UserInstallation=file://{profile_dir}",
                "--convert-to", output_fmt,
                "--outdir", tmpdir,
                input_path,
            ],
            capture_output=True,
            timeout=120,
        )
        if result.returncode != 0:
            logger.error(f"LibreOffice conversion failed: {result.stderr.decode()}")
            raise RuntimeError(f"LibreOffice conversion failed: {result.stderr.decode()}")

        output_path = os.path.join(tmpdir, f"input.{output_fmt}")
        if not os.path.exists(output_path):
            raise RuntimeError(f"LibreOffice did not produce a .{output_fmt} output file")

        with open(output_path, "rb") as f:
            return f.read()


def presentation_to_pdf(file_bytes, fmt):
    """Convert a presentation (PPT or PPTX) to PDF bytes using LibreOffice."""
    if fmt == "pptx":
        file_bytes = strip_presentation_backgrounds(file_bytes)
    return _libreoffice_convert(file_bytes, fmt, "pdf")


def convert_ppt_to_pptx(file_bytes):
    """Convert PPT bytes to PPTX bytes using LibreOffice."""
    return _libreoffice_convert(file_bytes, "ppt", "pptx")


async def extract_pptx_to_markdown(file_bytes):
    prs = Presentation(BytesIO(file_bytes))
    total_slides = len(prs.slides)

    all_parts = []
    all_image_tasks = []
    for slide_idx, slide in enumerate(prs.slides):
        slide_md, image_tasks = extract_slide(slide, slide_idx)
        all_parts.append((slide_idx, slide_md))
        all_image_tasks.extend(image_tasks)

    # Describe all images in parallel
    descriptions = await describe_images_batch(all_image_tasks)

    # Assemble final markdown
    markdown_parts = []
    for slide_idx, slide_md in all_parts:
        slide_md = resolve_image_placeholders(slide_md, descriptions)
        if slide_md.strip():
            markdown_parts.append(f"# Slide {slide_idx + 1}\n\n{slide_md}")

    combined = "\n\n---\n\n".join(markdown_parts) if markdown_parts else "# Conversion Failed\n\nNo content could be extracted."

    return ConvertResponse(
        markdown=combined,
        total_pages=total_slides,
        successful_pages=len(markdown_parts),
        failed_pages=total_slides - len(markdown_parts)
    )


@app.post("/convert", response_model=ConvertResponse)
async def convert_pdf(request: ConvertRequest):
    pdf_bytes = base64.b64decode(request.pdf_data)

    if len(pdf_bytes) / (1024 * 1024) > MAX_FILE_SIZE_MB:
        raise HTTPException(400, f"PDF exceeds {MAX_FILE_SIZE_MB}MB")

    chunks = render_pdf_to_chunks(pdf_bytes)
    return await process_chunks(chunks)


@app.post("/convert-pptx", response_model=ConvertResponse)
async def convert_pptx(request: ConvertPresentationRequest):
    file_bytes = base64.b64decode(request.presentation_data)

    if len(file_bytes) / (1024 * 1024) > MAX_FILE_SIZE_MB:
        raise HTTPException(400, f"File exceeds {MAX_FILE_SIZE_MB}MB")

    if request.format == "ppt":
        try:
            file_bytes = convert_ppt_to_pptx(file_bytes)
        except Exception as e:
            raise HTTPException(500, f"Failed to convert PPT to PPTX: {e}")

    return await extract_pptx_to_markdown(file_bytes)


@app.post("/convert-presentation-to-pdf", response_model=ConvertPresentationToPdfResponse)
async def convert_presentation_to_pdf(request: ConvertPresentationRequest):
    file_bytes = base64.b64decode(request.presentation_data)

    if len(file_bytes) / (1024 * 1024) > MAX_FILE_SIZE_MB:
        raise HTTPException(400, f"File exceeds {MAX_FILE_SIZE_MB}MB")

    try:
        pdf_bytes = presentation_to_pdf(file_bytes, request.format)
    except Exception as e:
        raise HTTPException(500, f"Presentation to PDF conversion failed: {e}")

    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    page_count = len(doc)
    doc.close()

    return ConvertPresentationToPdfResponse(
        pdf_data=base64.b64encode(pdf_bytes).decode(),
        pages=page_count,
    )


@app.get("/health")
async def health():
    return {"status": "ok"}
